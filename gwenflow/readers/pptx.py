import io
from dataclasses import dataclass
from pathlib import Path
from typing import List, Union

from gwenflow.logger import logger
from gwenflow.readers.base import Reader
from gwenflow.types import Document


@dataclass
class PptxReader(Reader):
    def __post_init__(self) -> None:
        try:
            __import__("pptx")
        except ImportError:
            raise ImportError("Missing required package: python-pptx. Install with: `uv add python-pptx`")

    def read(self, file: Union[Path, io.BytesIO]) -> List[Document]:
        from pptx import Presentation

        try:
            filename = self.get_file_name(file)
            content = self.get_file_content(file)
            prs = Presentation(content)
            documents = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                documents.append(
                    Document(
                        id=self.key(f"{filename}_slide{slide_num}"),
                        content="\n".join(texts),
                        metadata={"filename": filename, "page": slide_num, "slide": slide_num},
                    )
                )
            return documents
        except Exception as e:
            logger.exception(f"Error reading file: {e}")
            return []
